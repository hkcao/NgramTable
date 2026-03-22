#!/usr/bin/env python3
"""
N-gram Speculative Decoding 技术洞察与分析 PPT 生成脚本
基于 ngram_proposer.py 实际代码实现
清爽浅色风格
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Color Palette (清爽浅色系) ──
BG_WHITE  = RGBColor(0xFA, 0xFA, 0xFC)   # 近白背景
BG_CARD   = RGBColor(0xF0, 0xF4, 0xF8)   # 浅灰蓝卡片
BG_CODE   = RGBColor(0xF6, 0xF8, 0xFA)   # 代码块浅灰
BLUE      = RGBColor(0x22, 0x6C, 0xE0)   # 主色: 蓝
BLUE_L    = RGBColor(0xDB, 0xE8, 0xFD)   # 浅蓝背景
RED       = RGBColor(0xE0, 0x3E, 0x3E)   # 红
RED_L     = RGBColor(0xFD, 0xE8, 0xE8)   # 浅红
GREEN     = RGBColor(0x1A, 0x8C, 0x5B)   # 绿
GREEN_L   = RGBColor(0xDE, 0xF5, 0xE9)   # 浅绿
ORANGE    = RGBColor(0xE8, 0x8D, 0x14)   # 橙
ORANGE_L  = RGBColor(0xFE, 0xF3, 0xDB)   # 浅橙
PURPLE    = RGBColor(0x7C, 0x3A, 0xED)   # 紫
PURPLE_L  = RGBColor(0xED, 0xE5, 0xFD)   # 浅紫
BLACK     = RGBColor(0x1A, 0x1A, 0x2E)   # 深色文字
DARK      = RGBColor(0x33, 0x3D, 0x4C)   # 正文深灰
GRAY      = RGBColor(0x6B, 0x72, 0x80)   # 次要文字
LGRAY     = RGBColor(0x9C, 0xA3, 0xAF)   # 更淡灰
CODE_TEXT = RGBColor(0x24, 0x29, 0x2E)   # 代码文字

# ── Slide dimensions (16:9) ──
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

FONT = "Microsoft YaHei"
FONT_CODE = "Consolas"


# ══════════════════════════════════════════════════════════════
# Helper functions
# ══════════════════════════════════════════════════════════════

def set_slide_bg(slide, color=BG_WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text="",
             font_size=18, color=DARK, bold=False,
             alignment=PP_ALIGN.LEFT, font_name=FONT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_rich(slide, left, top, width, height, lines, font_name=FONT):
    """lines: list of (text, font_size, color, bold, alignment)"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, fs, clr, bld, align) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(fs)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = font_name
        p.alignment = align
        p.space_after = Pt(3)
    return tf


def add_card(slide, left, top, width, height,
             fill_color=BG_CARD, border_color=None, border_width=1.2):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()
    return shape


def add_code(slide, left, top, width, height, code_text, font_size=11):
    add_card(slide, left, top, width, height,
             fill_color=BG_CODE,
             border_color=RGBColor(0xD0, 0xD5, 0xDD))
    tf = add_text(slide, left + Inches(0.15), top + Inches(0.08),
                  width - Inches(0.3), height - Inches(0.16),
                  code_text, font_size=font_size,
                  color=CODE_TEXT, font_name=FONT_CODE)
    return tf


def add_flow_box(slide, left, top, width, height, text,
                 fill_color=BG_CARD, border_color=BLUE,
                 font_size=11, font_color=DARK):
    shape = add_card(slide, left, top, width, height,
                     fill_color=fill_color, border_color=border_color)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(font_size)
    tf.paragraphs[0].font.color.rgb = font_color
    tf.paragraphs[0].font.name = FONT
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape


def add_arrow(slide, left, top, width, height, color=BLUE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_down_arrow(slide, left, top, width, height, color=BLUE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_line(slide, left, top, width, color=BLUE, thickness=2.5):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(thickness))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_table(slide, left, top, rows, cols, col_widths, data,
              header_bg=BLUE, font_size=11):
    ts = slide.shapes.add_table(rows, cols, left, top,
                                 sum(col_widths), Inches(0.38) * rows)
    table = ts.table
    for ci, w in enumerate(col_widths):
        table.columns[ci].width = w
    for ri in range(rows):
        for ci in range(cols):
            cell = table.cell(ri, ci)
            cell.text = data[ri][ci]
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.name = FONT
                p.alignment = PP_ALIGN.CENTER
                if ri == 0:
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                else:
                    p.font.color.rgb = DARK
            cf = cell.fill
            cf.solid()
            if ri == 0:
                cf.fore_color.rgb = header_bg
            elif ri % 2 == 0:
                cf.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
            else:
                cf.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return table


def title_bar(slide, text, color=BLUE):
    """Page title with accent bar."""
    add_line(slide, Inches(0.5), Inches(0.6), Inches(0.6), color=color, thickness=4)
    add_text(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             text, font_size=28, color=BLACK, bold=True)


def flow_row(slide, y, steps, box_w=Inches(2.2), box_h=Inches(0.95),
             gap=Inches(0.35)):
    """Draw a horizontal flow of boxes with arrows."""
    for i, (text, clr) in enumerate(steps):
        x = Inches(0.5) + i * (box_w + gap)
        add_flow_box(slide, x, y, box_w, box_h, text, border_color=clr)
        if i < len(steps) - 1:
            add_arrow(slide, x + box_w + Inches(0.02),
                      y + Inches(0.28), Inches(0.3), Inches(0.35),
                      color=clr)


# ══════════════════════════════════════════════════════════════
# Slide 1: Title
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.0),
         "N-gram Speculative Decoding",
         font_size=42, color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)
add_line(slide, Inches(4.5), Inches(3.2), Inches(4.3), color=BLUE, thickness=3)
add_text(slide, Inches(1), Inches(3.4), Inches(11), Inches(0.7),
         "KMP  /  Hash  /  Trie (LookaheadCache)  /  Suffix Tree",
         font_size=22, color=BLUE, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.6),
         "技术实现原理  ·  关键设计  ·  对比分析",
         font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
         "基于 vLLM ngram_proposer.py + arctic-inference suffix_tree.cc 代码实现",
         font_size=13, color=LGRAY, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# Slide 2: Speculative Decoding Overview
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
title_bar(slide, "Speculative Decoding 概览")

# 核心思路
add_card(slide, Inches(0.5), Inches(1.0), Inches(12.3), Inches(1.2),
         fill_color=BLUE_L, border_color=BLUE)
add_rich(slide, Inches(0.8), Inches(1.05), Inches(11.7), Inches(1.1), [
    ("核心思路: 用廉价方法 \"猜\" 接下来的 k 个 token, 再用目标模型一次性验证", 15, BLUE, True, PP_ALIGN.LEFT),
    ("验证通过直接接受, 未通过则从拒绝点重新生成. 将 k 步自回归压缩为 1 步并行验证.", 13, DARK, False, PP_ALIGN.LEFT),
])

# Flow
y_flow = Inches(2.5)
boxes = [
    ("Context\n[t1, t2, ..., tn]", BLUE),
    ("N-gram Proposer\n猜 k 个 draft tokens", ORANGE),
    ("Target Model\n一次前向验证 k tokens", RED),
    ("Accept/Reject\n接受匹配, 拒绝后重新生成", GREEN),
]
box_w = Inches(2.7)
box_h = Inches(1.2)
gap = Inches(0.55)
for i, (text, clr) in enumerate(boxes):
    x = Inches(0.5) + i * (box_w + gap)
    fill = {BLUE: BLUE_L, ORANGE: ORANGE_L, RED: RED_L, GREEN: GREEN_L}[clr]
    add_flow_box(slide, x, y_flow, box_w, box_h, text,
                 fill_color=fill, border_color=clr, font_size=13, font_color=DARK)
    if i < len(boxes) - 1:
        add_arrow(slide, x + box_w + Inches(0.05), y_flow + Inches(0.4),
                  Inches(0.45), Inches(0.35), color=clr)

# 指标说明
add_card(slide, Inches(0.5), Inches(4.1), Inches(12.3), Inches(3.0),
         fill_color=BG_CARD)
add_rich(slide, Inches(0.8), Inches(4.2), Inches(11.7), Inches(2.8), [
    ("N-gram 方案: 不需要额外模型, 仅从已有 token 序列中挖掘重复模式来预测", 15, BLUE, True, PP_ALIGN.LEFT),
    ("", 6, DARK, False, PP_ALIGN.LEFT),
    ("Speedup = 原始耗时 / 投机解码耗时  (越高越好)", 13, DARK, False, PP_ALIGN.LEFT),
    ("Accept Rate = 被模型验证通过的 draft tokens / 总 draft tokens  (越高越好)", 13, DARK, False, PP_ALIGN.LEFT),
    ("Mean Draft Length = 平均每步生成的 draft token 数  (并非越多越好, 质量 > 数量)", 13, DARK, False, PP_ALIGN.LEFT),
    ("", 8, DARK, False, PP_ALIGN.LEFT),
    ("四种 N-gram Proposer:", 14, ORANGE, True, PP_ALIGN.LEFT),
    ("KMP (基线, 无状态)  ->  Hash (频率统计)  ->  Trie (树结构)  ->  Suffix Tree (C++ 极致)", 14, DARK, False, PP_ALIGN.LEFT),
])


# ══════════════════════════════════════════════════════════════
# Slide 3: KMP 设计逻辑
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
title_bar(slide, "1. KMP 方案 — vLLM 原生基线", color=BLUE)

# 左: 设计思路
add_card(slide, Inches(0.5), Inches(1.0), Inches(5.8), Inches(2.6),
         fill_color=BLUE_L, border_color=BLUE)
add_rich(slide, Inches(0.7), Inches(1.05), Inches(5.4), Inches(2.5), [
    ("设计逻辑", 17, BLUE, True, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("  不维护任何持久数据结构, 每次 propose 实时搜索", 13, DARK, False, PP_ALIGN.LEFT),
    ("  在当前序列中寻找与末尾匹配的最长 n-gram", 13, DARK, False, PP_ALIGN.LEFT),
    ("  取匹配位置之后的 k 个 token 作为 draft", 13, DARK, False, PP_ALIGN.LEFT),
    ("  Numba JIT 编译 + prange 批量并行加速", 13, DARK, False, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("  时间 O(n) | 空间 O(max_ngram) | 无预处理", 12, GRAY, False, PP_ALIGN.LEFT),
])

# 右: 关键技术
add_card(slide, Inches(6.8), Inches(1.0), Inches(6.0), Inches(2.6),
         fill_color=ORANGE_L, border_color=ORANGE)
add_rich(slide, Inches(7.0), Inches(1.05), Inches(5.6), Inches(2.5), [
    ("关键技术: 反转序列 + LPS 数组", 17, ORANGE, True, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("  1. 将 token 序列反转 (suffix -> prefix 问题转换)", 13, DARK, False, PP_ALIGN.LEFT),
    ("  2. 构建 LPS 数组 (长度限制为 max_ngram)", 13, DARK, False, PP_ALIGN.LEFT),
    ("     lps[i] = max{v : tokens[0:v] == tokens[i+1-v:i+1]}", 12, GRAY, False, PP_ALIGN.LEFT),
    ("  3. 扫描时追踪最长 prefix=suffix 匹配", 13, DARK, False, PP_ALIGN.LEFT),
    ("  4. 到达 max_ngram 时回退: prev_lps = lps[max_n-1]", 13, DARK, False, PP_ALIGN.LEFT),
    ("  5. 匹配位置反转回原序列, 取后续 k token", 13, DARK, False, PP_ALIGN.LEFT),
])

# 流程图
add_text(slide, Inches(0.5), Inches(3.8), Inches(3), Inches(0.4),
         "执行流程", font_size=15, color=BLUE, bold=True)
flow_row(slide, Inches(4.2), [
    ("反转序列\ntokens[::-1]", BLUE),
    ("构建 LPS\n(max_ngram 长度)", ORANGE),
    ("KMP 扫描\n追踪 longest_ngram", GREEN),
    ("位置反转\nstart=N-1-pos+len", RED),
    ("提取 draft\ntokens[start:start+k]", PURPLE),
])

# 特点总结
add_card(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.6),
         fill_color=BG_CARD)
add_rich(slide, Inches(0.8), Inches(5.55), Inches(11.7), Inches(1.5), [
    ("特点与局限", 15, BLUE, True, PP_ALIGN.LEFT),
    ("+ 零存储开销, 无需预处理          + Numba JIT O(n)          + prange 批量并行", 12, GREEN, False, PP_ALIGN.LEFT),
    ("- 每次从头扫描全序列 (重复计算)    - 只返回单条贪心链 (无分支)    - 长序列下变慢", 12, RED, False, PP_ALIGN.LEFT),
    ("Benchmark:  Speedup 1.70x  |  Accept 50.6%  |  MeanLen 3.52", 13, ORANGE, True, PP_ALIGN.LEFT),
])


# ══════════════════════════════════════════════════════════════
# Slide 4: KMP 核心代码
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
title_bar(slide, "1. KMP — 核心算法代码 (Numba JIT)", color=BLUE)

code_kmp = """@jit(nopython=True)
def _find_longest_matched_ngram_and_propose_tokens(
    origin_tokens, min_ngram, max_ngram, max_model_len, k):

    tokens = origin_tokens[::-1]               # Step 1: 反转序列
    lps = np.zeros(max_ngram, dtype=np.int32)  # Step 2: LPS 数组 (仅 max_ngram 长)

    longest_ngram = 0
    position = 0
    prev_lps = 0
    i = 1
    while i < total_token:                     # Step 3: KMP 线性扫描
        if tokens[prev_lps] == tokens[i]:
            prev_lps += 1
            if prev_lps >= longest_ngram:      #   更新最长匹配 (>=: 选最早出现)
                longest_ngram = prev_lps
                position = i
            if i < max_ngram:
                lps[i] = prev_lps              #   记录 LPS
            if prev_lps == max_ngram:          #   到达上限, 回退避免越界
                prev_lps = lps[max_ngram - 1]
            i += 1
        elif prev_lps != 0:
            prev_lps = lps[prev_lps - 1]      #   KMP 失配回退
        else:
            i += 1

    if longest_ngram < min_ngram:
        return np.empty((0,), dtype=origin_tokens.dtype)

    # Step 4: 反转位置, 提取 draft tokens
    start_position = total_token - 1 - position + longest_ngram
    return origin_tokens[start_position : start_position + k]"""

add_code(slide, Inches(0.5), Inches(0.9), Inches(8.8), Inches(5.8),
         code_kmp, font_size=12)

add_card(slide, Inches(9.6), Inches(0.9), Inches(3.3), Inches(5.8),
         fill_color=BLUE_L, border_color=BLUE)
add_rich(slide, Inches(9.8), Inches(1.0), Inches(2.9), Inches(5.6), [
    ("算法要点", 16, BLUE, True, PP_ALIGN.LEFT),
    ("", 6, DARK, False, PP_ALIGN.LEFT),
    ("反转的巧妙之处", 13, ORANGE, True, PP_ALIGN.LEFT),
    ("将 \"找末尾匹配的子串\"", 12, DARK, False, PP_ALIGN.LEFT),
    ("转化为 \"找开头匹配的", 12, DARK, False, PP_ALIGN.LEFT),
    ("前缀\" — KMP 经典问题", 12, DARK, False, PP_ALIGN.LEFT),
    ("", 6, DARK, False, PP_ALIGN.LEFT),
    ("LPS 截断", 13, ORANGE, True, PP_ALIGN.LEFT),
    ("只分配 max_ngram 大小", 12, DARK, False, PP_ALIGN.LEFT),
    ("达到上限时 fallback 到", 12, DARK, False, PP_ALIGN.LEFT),
    ("lps[max_n-1], 避免越界", 12, DARK, False, PP_ALIGN.LEFT),
    ("", 6, DARK, False, PP_ALIGN.LEFT),
    ("position 选择 (>=)", 13, ORANGE, True, PP_ALIGN.LEFT),
    ("倾向选最早出现的匹配", 12, DARK, False, PP_ALIGN.LEFT),
    ("(反转后 = 最靠后位置,", 12, DARK, False, PP_ALIGN.LEFT),
    (" 原序列中 = 最靠前)", 12, DARK, False, PP_ALIGN.LEFT),
    ("", 6, DARK, False, PP_ALIGN.LEFT),
    ("批量并行", 13, ORANGE, True, PP_ALIGN.LEFT),
    ("@njit(parallel=True)", 12, BLUE, False, PP_ALIGN.LEFT),
    ("prange 多请求并发", 12, DARK, False, PP_ALIGN.LEFT),
])


# ══════════════════════════════════════════════════════════════
# Slide 5: Hash 设计逻辑
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
title_bar(slide, "2. Hash 频率表方案", color=GREEN)

# 数据结构
add_card(slide, Inches(0.5), Inches(1.0), Inches(6.0), Inches(2.8),
         fill_color=GREEN_L, border_color=GREEN)
add_rich(slide, Inches(0.7), Inches(1.05), Inches(5.6), Inches(2.7), [
    ("数据结构 — 三层频率表", 17, GREEN, True, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("FreqTable: dict[tuple, Counter]", 13, BLUE, False, PP_ALIGN.LEFT),
    ("    (t1,t2,...tn) -> {next_token: count}", 12, GRAY, False, PP_ALIGN.LEFT),
    ("HashTable: dict[int, int]", 13, BLUE, False, PP_ALIGN.LEFT),
    ("    hash(context) -> best_token (快速路径)", 12, GRAY, False, PP_ALIGN.LEFT),
    ("LocalFreq: dict[req_id, dict[tuple, Counter]]", 13, BLUE, False, PP_ALIGN.LEFT),
    ("    每请求独立, 权重 3.0x (近因效应)", 12, GRAY, False, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("合并: merged = global + local * 3.0", 13, GREEN, True, PP_ALIGN.LEFT),
])

# 设计思路
add_card(slide, Inches(6.8), Inches(1.0), Inches(6.0), Inches(2.8),
         fill_color=ORANGE_L, border_color=ORANGE)
add_rich(slide, Inches(7.0), Inches(1.05), Inches(5.6), Inches(2.7), [
    ("设计逻辑", 17, ORANGE, True, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("1. 滑动窗口收集 n-gram 频率: n in [min_n, max_n]", 13, DARK, False, PP_ALIGN.LEFT),
    ("2. 增量更新: 只处理包含新 token 的窗口", 13, DARK, False, PP_ALIGN.LEFT),
    ("    start = max(0, total - n_new - n)", 12, GRAY, False, PP_ALIGN.LEFT),
    ("3. 查询: 从长 n-gram 到短, 取最高置信度", 13, DARK, False, PP_ALIGN.LEFT),
    ("4. 三重过滤: 置信度 + 投票 + 自回归链", 13, DARK, False, PP_ALIGN.LEFT),
    ("5. 每 100 次更新异步持久化到磁盘", 13, DARK, False, PP_ALIGN.LEFT),
])

# Propose 流程
add_text(slide, Inches(0.5), Inches(4.0), Inches(5), Inches(0.4),
         "Propose 流程 (每步迭代 k 次)", font_size=15, color=GREEN, bold=True)
flow_row(slide, Inches(4.4), [
    ("长 n-gram 优先\ntry_n: max -> min", GREEN),
    ("合并查询\nglobal + local * 3.0", ORANGE),
    ("置信度检查\ntop/total >= 0.3?", RED),
    ("多 n-gram 投票\n多数窗口同意?", PURPLE),
    ("自回归扩展\nextended.append(tok)", BLUE),
])

# 三重质量控制
add_card(slide, Inches(0.5), Inches(5.7), Inches(12.3), Inches(1.5),
         fill_color=BG_CARD)
add_rich(slide, Inches(0.8), Inches(5.75), Inches(11.7), Inches(1.4), [
    ("三重质量控制", 15, GREEN, True, PP_ALIGN.LEFT),
    ("1) 置信度过滤: top_count/total < 0.3 时停止 — 不确定就不猜", 12, DARK, False, PP_ALIGN.LEFT),
    ("2) 多 n-gram 投票: [min_n..max_n] 各窗口半数以上同意才继续 — 交叉验证", 12, DARK, False, PP_ALIGN.LEFT),
    ("3) 请求局部叠加: 当前请求 pattern 权重 3x, 新请求时重置 — 突出当前上下文", 12, DARK, False, PP_ALIGN.LEFT),
    ("Benchmark:  Speedup 1.86x  |  Accept 43.3%  |  MeanLen 3.03", 13, ORANGE, True, PP_ALIGN.LEFT),
])


# ══════════════════════════════════════════════════════════════
# Slide 6: Hash 核心代码
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
title_bar(slide, "2. Hash — 核心 Propose + 投票代码", color=GREEN)

code_hash = """def _propose_tokens_hash(self, input_ids, k, req_idx=0):
    extended = list(input_ids)
    for step in range(k):
        best_token, best_confidence = None, 0.0
        for try_n in range(n, self.min_n - 1, -1):   # 长 n-gram 优先
            context = tuple(extended[-try_n:])
            counter = self._get_merged_counter(context, req_idx)
            #  ^ 关键: global_freq + local_freq * 3.0
            if counter:
                token, conf = self._check_confidence(counter)
                #  ^ conf = top_count / total_count
                if token and conf > best_confidence:
                    best_token, best_confidence = token, conf
                    if conf >= self._min_confidence:   # 0.3
                        break
        if best_token is None:           break
        if best_confidence < 0.3:        break  # 置信度不足
        if not self._vote_ngrams(extended, best_token, req_idx):
            break                                # 投票不通过
        drafts.append(best_token)
        extended.append(best_token)              # 自回归扩展"""

add_code(slide, Inches(0.5), Inches(0.9), Inches(8.0), Inches(4.2),
         code_hash, font_size=12)

code_vote = """def _vote_ngrams(self, input_ids, candidate, req_idx):
    votes_for = votes_total = 0
    for n in range(self.min_n, n_max + 1):
        context = tuple(input_ids[-n:])
        counter = self._get_merged_counter(context, req_idx)
        if counter:
            votes_total += 1
            best, _ = self._check_confidence(counter)
            if best == candidate:
                votes_for += 1
    return votes_for >= (votes_total + 1) // 2  # 多数同意"""

add_code(slide, Inches(0.5), Inches(5.3), Inches(8.0), Inches(2.0),
         code_vote, font_size=12)

# 右侧说明
add_card(slide, Inches(8.8), Inches(0.9), Inches(4.1), Inches(3.0),
         fill_color=GREEN_L, border_color=GREEN)
add_rich(slide, Inches(9.0), Inches(1.0), Inches(3.7), Inches(2.8), [
    ("频率合并策略", 15, GREEN, True, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("merged = Counter(global)", 12, BLUE, False, PP_ALIGN.LEFT),
    ("for tok, cnt in local:", 12, BLUE, False, PP_ALIGN.LEFT),
    ("  merged[tok] += cnt * 3.0", 12, BLUE, False, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("local 在新请求时清空重建", 12, DARK, False, PP_ALIGN.LEFT),
    ("对当前对话的重复 pattern", 12, DARK, False, PP_ALIGN.LEFT),
    ("给予更高权重 (近因偏好)", 12, DARK, False, PP_ALIGN.LEFT),
])

add_card(slide, Inches(8.8), Inches(4.2), Inches(4.1), Inches(3.1),
         fill_color=ORANGE_L, border_color=ORANGE)
add_rich(slide, Inches(9.0), Inches(4.3), Inches(3.7), Inches(2.9), [
    ("增量更新机制", 15, ORANGE, True, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("新请求: 全量构建 FreqTable", 12, DARK, False, PP_ALIGN.LEFT),
    ("续写: 仅更新包含新 token", 12, DARK, False, PP_ALIGN.LEFT),
    ("的滑动窗口, 并检查:", 12, DARK, False, PP_ALIGN.LEFT),
    ("  if new_best != old_best:", 12, BLUE, False, PP_ALIGN.LEFT),
    ("    hash_table[h] = new", 12, BLUE, False, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("每 100 次更新异步 flush", 12, DARK, False, PP_ALIGN.LEFT),
    ("到磁盘 (pickle 序列化)", 12, DARK, False, PP_ALIGN.LEFT),
])


# ══════════════════════════════════════════════════════════════
# Slide 7: Trie 设计逻辑
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
title_bar(slide, "3. Trie 方案 — LookaheadCache 复现", color=PURPLE)

# 三层结构
add_card(slide, Inches(0.5), Inches(1.0), Inches(4.0), Inches(3.2),
         fill_color=PURPLE_L, border_color=PURPLE)
add_rich(slide, Inches(0.7), Inches(1.05), Inches(3.6), Inches(3.1), [
    ("三层数据结构", 16, PURPLE, True, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("TrieCache (顶层)", 14, BLUE, True, PP_ALIGN.LEFT),
    ("  mem: root_key -> TrieTree", 12, DARK, False, PP_ALIGN.LEFT),
    ("  root_key = 1-token 或 n-tuple", 12, GRAY, False, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("TrieTree (每个根一棵)", 14, BLUE, True, PP_ALIGN.LEFT),
    ("  nodes: token_id -> TrieNode", 12, DARK, False, PP_ALIGN.LEFT),
    ("  max 65536 节点, 超限剪枝", 12, GRAY, False, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("TrieNode", 14, BLUE, True, PP_ALIGN.LEFT),
    ("  freqs: {idx: float}  频率", 12, DARK, False, PP_ALIGN.LEFT),
    ("  children: {token: TrieNode}", 12, DARK, False, PP_ALIGN.LEFT),
])

# 频率语义
add_card(slide, Inches(4.8), Inches(1.0), Inches(4.0), Inches(3.2),
         fill_color=GREEN_L, border_color=GREEN)
add_rich(slide, Inches(5.0), Inches(1.05), Inches(3.6), Inches(3.1), [
    ("频率双通道", 16, GREEN, True, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("freqs[idx] = 输入频率", 14, DARK, True, PP_ALIGN.LEFT),
    ("  来自 prompt (input mode)", 12, GRAY, False, PP_ALIGN.LEFT),
    ("  每个 request 独立 idx", 12, GRAY, False, PP_ALIGN.LEFT),
    ("  请求结束后重置为 0", 12, GRAY, False, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("freqs[-1] = 输出频率", 14, DARK, True, PP_ALIGN.LEFT),
    ("  来自生成 (output mode)", 12, GRAY, False, PP_ALIGN.LEFT),
    ("  跨请求共享, 持久积累", 12, GRAY, False, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("mix 模式混合:", 13, GREEN, True, PP_ALIGN.LEFT),
    ("fm = (1-w)*fi + w*fo  (w=0.5)", 12, BLUE, False, PP_ALIGN.LEFT),
])

# Root key
add_card(slide, Inches(9.1), Inches(1.0), Inches(3.7), Inches(3.2),
         fill_color=ORANGE_L, border_color=ORANGE)
add_rich(slide, Inches(9.3), Inches(1.05), Inches(3.3), Inches(3.1), [
    ("Root Key 设计", 16, ORANGE, True, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("node_size=1 (默认)", 13, DARK, True, PP_ALIGN.LEFT),
    ("  root = single token", 12, GRAY, False, PP_ALIGN.LEFT),
    ("  碰撞多, 树大而浅", 12, GRAY, False, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("node_size=3 (最优)", 13, GREEN, True, PP_ALIGN.LEFT),
    ("  root = (t1, t2, t3)", 12, DARK, False, PP_ALIGN.LEFT),
    ("  精确定位, 内部仍单 token", 12, DARK, False, PP_ALIGN.LEFT),
    ("  本质 ≡ suffix min_match", 12, GRAY, False, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("_make_key 始终返回", 13, DARK, True, PP_ALIGN.LEFT),
    ("  (single_token, 1)", 12, BLUE, False, PP_ALIGN.LEFT),
])

# 写入 / 查询
add_text(slide, Inches(0.5), Inches(4.4), Inches(5), Inches(0.3),
         "写入: 滑动窗口 + 链式分支插入", font_size=14, color=PURPLE, bold=True)
flow_row(slide, Inches(4.75), [
    ("滑动窗口\ni: 0 -> len-ns", PURPLE),
    ("提取 root_key\n+ branch tokens", ORANGE),
    ("链式插入 _put\n沿路径递增频率", GREEN),
    ("流式增量\nstream_put", BLUE),
], box_w=Inches(2.7), gap=Inches(0.5))

add_text(slide, Inches(0.5), Inches(5.9), Inches(5), Inches(0.3),
         "查询: hier_get -> get -> 提取最优分支",
         font_size=14, color=PURPLE, bold=True)
flow_row(slide, Inches(6.25), [
    ("滑窗找匹配\nroot_key -> TrieTree", PURPLE),
    ("_match 精确匹配\n沿 Trie 下行", ORANGE),
    ("_ravel DFS\n频率降序展开子树", GREEN),
    ("提取贪心路径\nmask[j,cur]==1", BLUE),
], box_w=Inches(2.7), gap=Inches(0.5))


# ══════════════════════════════════════════════════════════════
# Slide 8: Trie 核心代码 — 写入与匹配
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
title_bar(slide, "3. Trie — 写入 + 匹配核心代码", color=PURPLE)

code_trie_put = """# TrieCache.put — 滑动窗口插入
def put(self, token_ids, branch_length=8, mode='output', idx=0):
    ns = self.node_size  # e.g. 3
    for i in range(len(token_ids) - ns + 1):
        root_key = token_ids[i] if ns == 1 \\
            else tuple(token_ids[i:i + ns])
        tup = token_ids[i + ns : i + ns + branch_length]
        tid = root_key if ns == 1 else root_key[-1]
        self._put_to_tree(root_key, tup, mode, idx, tid)

# TrieTree._put — 沿路径插入, 已有节点递增频率
def _put(self, token_ids, nodes, mode='output', freq=1.0, idx=-1):
    pos = 0
    while pos < len(token_ids):
        key, width = self._make_key(token_ids, pos)  # always (token, 1)
        node = nodes.get(key, None)
        if node is None:
            nodes.update(self._pack(token_ids, pos, idx, freq=freq))
            break                       # 一次性构建剩余链
        node.freqs[idx] = node.freqs.get(idx, 0.0) + freq  # 递增频率
        nodes = node.children
        pos += width"""

add_code(slide, Inches(0.5), Inches(0.9), Inches(8.0), Inches(3.7),
         code_trie_put, font_size=11)

code_trie_match = """# TrieTree._match — 精确路径匹配
def _match(self, token_ids, mode='mix', idx=0, skip_gram=False):
    nodes = self.nodes
    token_id = None
    pos = 0
    while pos < len(token_ids):
        key, width = self._make_key(token_ids, pos)
        node = nodes.get(key, None)
        if node is None and skip_gram:          # skip-gram 回退
            node = nodes.get(_SKIPGRAM_SENTINEL, None)
        if node is None:
            break                               # 匹配终止
        token_id = token_ids[pos + width - 1]
        # 频率检查: 该路径是否有有效频率
        if mode == 'mix':
            if node.freqs.get(idx, 0) > 0 or node.freqs.get(-1, 0) > 0:
                nodes = node.children           # 有效, 继续下行
            else:
                nodes = {}
        pos += width
    return token_id, nodes  # (最后匹配 token, 子节点集合)"""

add_code(slide, Inches(0.5), Inches(4.8), Inches(8.0), Inches(2.6),
         code_trie_match, font_size=11)

# 右侧注解
add_card(slide, Inches(8.8), Inches(0.9), Inches(4.1), Inches(3.2),
         fill_color=PURPLE_L, border_color=PURPLE)
add_rich(slide, Inches(9.0), Inches(1.0), Inches(3.7), Inches(3.0), [
    ("写入路径要点", 15, PURPLE, True, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("_pack: 一次性构建", 13, DARK, True, PP_ALIGN.LEFT),
    ("新节点从尾到头倒序链接", 12, DARK, False, PP_ALIGN.LEFT),
    ("避免多次 dict 插入", 12, GRAY, False, PP_ALIGN.LEFT),
    ("", 6, DARK, False, PP_ALIGN.LEFT),
    ("已有路径: 仅递增频率", 13, DARK, True, PP_ALIGN.LEFT),
    ("node.freqs[idx] += freq", 12, BLUE, False, PP_ALIGN.LEFT),
    ("不创建新节点, O(depth)", 12, GRAY, False, PP_ALIGN.LEFT),
    ("", 6, DARK, False, PP_ALIGN.LEFT),
    ("stream_put 流式增量", 13, DARK, True, PP_ALIGN.LEFT),
    ("缓冲生成 token, 攒够", 12, DARK, False, PP_ALIGN.LEFT),
    ("branch_length 再插入", 12, DARK, False, PP_ALIGN.LEFT),
])

add_card(slide, Inches(8.8), Inches(4.4), Inches(4.1), Inches(3.0),
         fill_color=GREEN_L, border_color=GREEN)
add_rich(slide, Inches(9.0), Inches(4.5), Inches(3.7), Inches(2.8), [
    ("匹配路径要点", 15, GREEN, True, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("精确逐 token 下行", 13, DARK, True, PP_ALIGN.LEFT),
    ("每步检查频率有效性", 12, DARK, False, PP_ALIGN.LEFT),
    ("无效则中断 (nodes={})", 12, GRAY, False, PP_ALIGN.LEFT),
    ("", 6, DARK, False, PP_ALIGN.LEFT),
    ("返回 (token, children)", 13, DARK, True, PP_ALIGN.LEFT),
    ("children 用于后续 DFS", 12, DARK, False, PP_ALIGN.LEFT),
    ("展开候选 draft tokens", 12, GRAY, False, PP_ALIGN.LEFT),
    ("", 6, DARK, False, PP_ALIGN.LEFT),
    ("skip_gram 回退", 13, DARK, True, PP_ALIGN.LEFT),
    ("精确匹配失败时尝试", 12, DARK, False, PP_ALIGN.LEFT),
    ("sentinel(-999) 子节点", 12, GRAY, False, PP_ALIGN.LEFT),
])


# ══════════════════════════════════════════════════════════════
# Slide 9: Trie 核心代码 — DFS 展开 + 最优路径提取
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
title_bar(slide, "3. Trie — DFS 展开 + 最优路径提取", color=PURPLE)

code_ravel = """# TrieTree._ravel — 多分支 DFS 按频率降序展开
def _ravel(self, nodes, ids, mask, pid, max_size=64, max_length=8,
           output_weight=0.5, mode='mix', idx=0, ...):
    if len(ids) >= max_size or max_length <= 0:
        return
    # 子节点按混合频率降序排列
    sorts = sorted([
        (k, v, (1-output_weight)*v.freqs.get(idx, 0)
               + output_weight*v.freqs.get(-1, 0))
        for k, v in nodes.items()
    ], key=lambda x: x[2], reverse=True)

    for tid, node, fm in sorts:
        if len(ids) >= max_size: return
        fi, fo = node.freqs.get(idx, 0), node.freqs.get(-1, 0)
        if fi < min_input_freq and fo < min_output_freq and fm < min_mix_freq:
            continue                    # 频率不够, 裁剪
        ids.append(tid)
        rid = len(ids) - 1
        if pid > -1: mask[rid] = mask[pid]  # 继承祖先关系
        mask[rid, rid] = 1
        if node.children:
            self._ravel(node.children, ids, mask, rid, ...)  # 递归"""

add_code(slide, Inches(0.5), Inches(0.9), Inches(8.0), Inches(4.0),
         code_ravel, font_size=11)

code_extract = """# TrieCache.get — 从多分支中提取贪心最优路径
def get(self, token_ids, ...):
    decoding_ids, decoding_masks, sizes = self.hier_get(token_ids, ...)
    # ids[0] 是 root (context), 不算 draft
    # _ravel 按频率降序 DFS -> 第一条连续链就是最优路径
    branch = []
    current = 0
    for j in range(1, len(decoding_ids)):
        if decoding_masks[j, current] == 1:  # j 是 current 的后代
            branch.append(decoding_ids[j])
            current = j
        else:
            break   # 分支切换 = 不连续, 停止
    return branch"""

add_code(slide, Inches(0.5), Inches(5.1), Inches(8.0), Inches(2.3),
         code_extract, font_size=11)

# 右侧注解
add_card(slide, Inches(8.8), Inches(0.9), Inches(4.1), Inches(6.5),
         fill_color=PURPLE_L, border_color=PURPLE)
add_rich(slide, Inches(9.0), Inches(1.0), Inches(3.7), Inches(6.3), [
    ("DFS 展开结构", 16, PURPLE, True, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("多分支 Trie:", 13, DARK, True, PP_ALIGN.LEFT),
    ("  t4[15] -> t5[8] -> t6[3]", 12, GREEN, False, PP_ALIGN.LEFT),
    ("           -> t7[5] -> t8[2]", 12, ORANGE, False, PP_ALIGN.LEFT),
    ("  t9[4]  -> t10[2]", 12, GRAY, False, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("DFS 输出:", 13, DARK, True, PP_ALIGN.LEFT),
    ("ids=[root,t4,t5,t6,t7,t8,t9,t10]", 11, BLUE, False, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("mask 记录祖先关系:", 13, DARK, True, PP_ALIGN.LEFT),
    ("mask[j] = mask[parent]", 12, BLUE, False, PP_ALIGN.LEFT),
    ("mask[j,j] = 1", 12, BLUE, False, PP_ALIGN.LEFT),
    ("", 6, DARK, False, PP_ALIGN.LEFT),
    ("最优路径提取:", 16, GREEN, True, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("因为是频率降序 DFS", 13, DARK, False, PP_ALIGN.LEFT),
    ("第一条连续链即最优:", 13, DARK, False, PP_ALIGN.LEFT),
    ("  [t4, t5, t6]", 13, GREEN, True, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("mask[j, current]==1 检查", 12, DARK, False, PP_ALIGN.LEFT),
    ("j 是否在 current 的子树中", 12, DARK, False, PP_ALIGN.LEFT),
    ("不连续立即 break", 12, GRAY, False, PP_ALIGN.LEFT),
    ("", 6, DARK, False, PP_ALIGN.LEFT),
    ("Benchmark:", 13, ORANGE, True, PP_ALIGN.LEFT),
    ("Speedup 1.95x | Accept 55.7%", 12, DARK, False, PP_ALIGN.LEFT),
    ("MeanLen 3.67", 12, DARK, False, PP_ALIGN.LEFT),
])


# ══════════════════════════════════════════════════════════════
# Slide 10: Suffix Tree 设计逻辑
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
title_bar(slide, "4. Suffix Tree 方案 — C++ arctic-inference", color=RED)

# 数据结构
add_card(slide, Inches(0.5), Inches(1.0), Inches(4.2), Inches(3.3),
         fill_color=RED_L, border_color=RED)
add_rich(slide, Inches(0.7), Inches(1.05), Inches(3.8), Inches(3.2), [
    ("数据结构: 压缩后缀树", 16, RED, True, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("struct Node {", 12, BLUE, False, PP_ALIGN.LEFT),
    ("  int64 count;     // 经过的后缀数", 11, DARK, False, PP_ALIGN.LEFT),
    ("  int token;       // 首 token ID", 11, DARK, False, PP_ALIGN.LEFT),
    ("  int length;      // 路径压缩长度", 11, DARK, False, PP_ALIGN.LEFT),
    ("  int ref_seq;     // 引用序列 ID", 11, DARK, False, PP_ALIGN.LEFT),
    ("  int ref_idx;     // 引用起始位置", 11, DARK, False, PP_ALIGN.LEFT),
    ("  Map children;    // 子节点", 11, DARK, False, PP_ALIGN.LEFT),
    ("  Node* head_child; // count 最大子", 11, DARK, False, PP_ALIGN.LEFT),
    ("  Node* siblings;  // 兄弟双向链表", 11, DARK, False, PP_ALIGN.LEFT),
    ("}", 12, BLUE, False, PP_ALIGN.LEFT),
])

# 核心机制
add_card(slide, Inches(4.9), Inches(1.0), Inches(4.0), Inches(3.3),
         fill_color=ORANGE_L, border_color=ORANGE)
add_rich(slide, Inches(5.1), Inches(1.05), Inches(3.6), Inches(3.2), [
    ("核心设计", 16, ORANGE, True, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("1. 零冗余存储", 14, DARK, True, PP_ALIGN.LEFT),
    ("   tokens 只存一份在 _seqs", 12, DARK, False, PP_ALIGN.LEFT),
    ("   节点通过 (ref_seq, ref_idx)", 12, DARK, False, PP_ALIGN.LEFT),
    ("   引用原始数据, 不复制", 12, GRAY, False, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("2. 增量构建", 14, DARK, True, PP_ALIGN.LEFT),
    ("   append(seq_id, token)", 12, DARK, False, PP_ALIGN.LEFT),
    ("   active_nodes 维护活跃后缀", 12, DARK, False, PP_ALIGN.LEFT),
    ("   每 token O(depth) 更新", 12, GRAY, False, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("3. 概率估计 (内建)", 14, DARK, True, PP_ALIGN.LEFT),
    ("   prob = child.count/parent.count", 12, BLUE, False, PP_ALIGN.LEFT),
    ("   无需额外 FreqTable", 12, GRAY, False, PP_ALIGN.LEFT),
])

# Speculate 模式
add_card(slide, Inches(9.1), Inches(1.0), Inches(3.7), Inches(3.3),
         fill_color=BLUE_L, border_color=BLUE)
add_rich(slide, Inches(9.3), Inches(1.05), Inches(3.3), Inches(3.2), [
    ("两种推测模式", 16, BLUE, True, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("Path 模式 (贪心)", 13, GREEN, True, PP_ALIGN.LEFT),
    ("  始终选 head_child", 12, DARK, False, PP_ALIGN.LEFT),
    ("  (count 最大子节点)", 12, GRAY, False, PP_ALIGN.LEFT),
    ("  O(1) 取最优, 单链输出", 12, DARK, False, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("Tree 模式 (多分支)", 13, PURPLE, True, PP_ALIGN.LEFT),
    ("  优先队列 max-heap", 12, DARK, False, PP_ALIGN.LEFT),
    ("  按累积概率展开多分支", 12, DARK, False, PP_ALIGN.LEFT),
    ("  prob < min_token_prob", 12, DARK, False, PP_ALIGN.LEFT),
    ("  时剪枝", 12, GRAY, False, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("vLLM 默认 Path 模式", 13, ORANGE, True, PP_ALIGN.LEFT),
])

# 生命周期
add_text(slide, Inches(0.5), Inches(4.5), Inches(3), Inches(0.3),
         "请求生命周期", font_size=14, color=RED, bold=True)
flow_row(slide, Inches(4.85), [
    ("start_request\nextend(prompt)", RED),
    ("每步生成\nappend(token)", ORANGE),
    ("match_context\n匹配末尾 context", GREEN),
    ("speculate_path\nhead_child 贪心", BLUE),
    ("stop_request\nremove + prune", PURPLE),
])

# 底部
add_card(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(1.1),
         fill_color=BG_CARD)
add_rich(slide, Inches(0.8), Inches(6.15), Inches(11.7), Inches(1.0), [
    ("+ C++ 零 Python 开销   + 路径压缩 O(n) 空间   + 兄弟链 O(1) 最优子节点   + 概率天然内建", 12, GREEN, False, PP_ALIGN.LEFT),
    ("Benchmark:  Speedup 2.04x  |  Accept 54.2%  |  MeanLen 2.50  (所有 N-gram 方案最优)", 13, ORANGE, True, PP_ALIGN.LEFT),
])


# ══════════════════════════════════════════════════════════════
# Slide 11: Suffix Tree 核心代码
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
title_bar(slide, "4. Suffix Tree — match_context + speculate_path", color=RED)

code_match = """// suffix_tree.cc — _match_context: 在后缀树中匹配上下文
pair<Node*, int> _match_context(span<int> context) {
    Node* node = root;
    int idx = 0;              // 当前节点内的偏移
    for (int token : context) {
        if (idx >= node->length) {
            // 当前节点走完, 查找子节点
            auto it = node->children.find(token);
            if (it == node->children.end())
                return {nullptr, -1};   // 无匹配
            node = it->second;
            idx = 0;                   // 进入子节点第 0 个位置
        }
        // 检查节点内 token 是否匹配
        if (_seqs[node->ref_seq][node->ref_idx + idx] != token)
            return {nullptr, -1};
        idx++;
    }
    return {node, idx};  // 返回匹配节点 + 偏移
}"""

add_code(slide, Inches(0.5), Inches(0.9), Inches(7.5), Inches(3.7),
         code_match, font_size=11)

code_spec = """// suffix_tree.cc — speculate: 枚举 match_len 取最优
Draft speculate(span<int> context, int max_spec, ...) {
    Draft best;
    for (int ml = 1; ml < context.size(); ml++) {
        auto [node, idx] = _match_context(context.last(ml));
        if (!node) break;
        int max_tok = min(max_spec, int(ml * factor));
        Draft d = _speculate_path(node, idx, max_tok);
        if (d.score >= best.score) best = d;  // >= 偏好更长
    }
    return best;
}
// _speculate_path: 贪心跟随 head_child
Draft _speculate_path(Node* node, int idx, int max_tokens) {
    float prob = 1.0f;  // float32 精度
    while (draft.size() < max_tokens && prob >= min_prob) {
        if (idx < node->length) {
            draft.push(ref_token[idx], prob); idx++;
        } else {
            Node* child = node->head_child;  // O(1) 取最优
            if (!child) break;
            prob *= (float)child->count / node->count;
            node = child; idx = 0;
        }
    }
}"""

add_code(slide, Inches(0.5), Inches(4.8), Inches(7.5), Inches(2.6),
         code_spec, font_size=11)

# 右侧注解
add_card(slide, Inches(8.3), Inches(0.9), Inches(4.5), Inches(3.3),
         fill_color=RED_L, border_color=RED)
add_rich(slide, Inches(8.5), Inches(1.0), Inches(4.1), Inches(3.1), [
    ("match_context 要点", 15, RED, True, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("路径压缩匹配:", 13, DARK, True, PP_ALIGN.LEFT),
    ("一个节点内存多个 token", 12, DARK, False, PP_ALIGN.LEFT),
    ("通过 ref_seq[ref_idx+offset]", 12, BLUE, False, PP_ALIGN.LEFT),
    ("直接访问, 无需逐节点跳转", 12, DARK, False, PP_ALIGN.LEFT),
    ("", 6, DARK, False, PP_ALIGN.LEFT),
    ("失配即停:", 13, DARK, True, PP_ALIGN.LEFT),
    ("节点内 mismatch 直接返回 null", 12, DARK, False, PP_ALIGN.LEFT),
    ("子节点不存在也返回 null", 12, DARK, False, PP_ALIGN.LEFT),
    ("精确匹配, 不做模糊回退", 12, GRAY, False, PP_ALIGN.LEFT),
])

add_card(slide, Inches(8.3), Inches(4.5), Inches(4.5), Inches(2.9),
         fill_color=GREEN_L, border_color=GREEN)
add_rich(slide, Inches(8.5), Inches(4.6), Inches(4.1), Inches(2.7), [
    ("speculate 要点", 15, GREEN, True, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("枚举所有 match_len:", 13, DARK, True, PP_ALIGN.LEFT),
    ("for ml = 1 to len(ctx)-1", 12, BLUE, False, PP_ALIGN.LEFT),
    ("取 score 最高的 (≥ 偏好长)", 12, DARK, False, PP_ALIGN.LEFT),
    ("", 6, DARK, False, PP_ALIGN.LEFT),
    ("head_child = O(1):", 13, DARK, True, PP_ALIGN.LEFT),
    ("兄弟链按 count 降序排列", 12, DARK, False, PP_ALIGN.LEFT),
    ("等 count 时后提升者在前", 12, DARK, False, PP_ALIGN.LEFT),
    ("", 6, DARK, False, PP_ALIGN.LEFT),
    ("累积概率衰减:", 13, DARK, True, PP_ALIGN.LEFT),
    ("prob *= child.count/parent.count", 12, BLUE, False, PP_ALIGN.LEFT),
    ("float32 精度, < 阈值即停", 12, GRAY, False, PP_ALIGN.LEFT),
])


# ══════════════════════════════════════════════════════════════
# Slide 12: 对比总结表格
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
title_bar(slide, "方案对比总结", color=BLUE)

data = [
    ["维度", "KMP", "Hash", "Trie (LookaheadCache)", "Suffix Tree (C++)"],
    ["数据结构", "无 (每次重算)", "dict[tuple, Counter]\n+ 请求局部表", "三层 Trie 树\nTrieCache/Tree/Node", "压缩后缀树\n路径压缩+兄弟链"],
    ["存储", "O(1) 临时", "O(|patterns|)", "O(nodes) max 65536", "O(n) 路径压缩"],
    ["模式收集", "无 (实时搜索)", "滑窗 [min_n,max_n]\n增量更新", "滑窗 + 分支插入\nstream_put 增量", "全后缀索引\nappend 增量"],
    ["Propose", "反转+LPS\n取后续 k token", "长->短 n-gram\n置信度+投票", "滑窗找 root_key\nDFS+贪心路径", "match_context\nhead_child 贪心"],
    ["质量控制", "无", "置信度 0.3\n多 n-gram 投票", "频率排序\n输入/输出混合", "prob 概率估计\nmin_token_prob"],
    ["实现", "Numba JIT", "Python", "Python", "C++"],
    ["Speedup", "1.70x", "1.86x", "1.95x", "2.04x (最优)"],
    ["Accept", "50.6%", "43.3%", "55.7%", "54.2%"],
    ["Avg Draft", "3.52", "3.03", "3.67", "2.50 (最少)"],
]

col_widths = [Inches(1.6), Inches(2.3), Inches(2.8), Inches(3.0), Inches(3.0)]
add_table(slide, Inches(0.15), Inches(0.9), len(data), len(data[0]),
          col_widths, data, font_size=10)


# ══════════════════════════════════════════════════════════════
# Slide 13: 关键洞察
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
title_bar(slide, "关键技术洞察", color=BLUE)

# Insight 1
add_card(slide, Inches(0.5), Inches(1.0), Inches(6.0), Inches(1.6),
         fill_color=GREEN_L, border_color=GREEN)
add_rich(slide, Inches(0.7), Inches(1.05), Inches(5.6), Inches(1.5), [
    ("\"猜得准\" > \"猜得多\"", 17, GREEN, True, PP_ALIGN.LEFT),
    ("", 3, DARK, False, PP_ALIGN.LEFT),
    ("Suffix: 最少 draft (2.50) 但最高 speedup (2.04x)", 13, DARK, False, PP_ALIGN.LEFT),
    ("Trie: 最多 draft (3.67) 但 speedup 反而更低 (1.95x)", 13, DARK, False, PP_ALIGN.LEFT),
    ("错误 draft 浪费验证计算, 精准少量优于盲目多猜", 12, GREEN, False, PP_ALIGN.LEFT),
])

# Insight 2
add_card(slide, Inches(6.8), Inches(1.0), Inches(6.0), Inches(1.6),
         fill_color=RED_L, border_color=RED)
add_rich(slide, Inches(7.0), Inches(1.05), Inches(5.6), Inches(1.5), [
    ("所有模糊匹配均失败", 17, RED, True, PP_ALIGN.LEFT),
    ("", 3, DARK, False, PP_ALIGN.LEFT),
    ("Fuzzy (43%) / SkipGram (31%) / EditDist (28%)", 13, DARK, False, PP_ALIGN.LEFT),
    ("放松条件 -> 频率分散 -> 预测质量下降", 13, DARK, False, PP_ALIGN.LEFT),
    ("N-gram 是 context-free 的, 模糊化引入噪声而非信号", 12, RED, False, PP_ALIGN.LEFT),
])

# Insight 3
add_card(slide, Inches(0.5), Inches(2.9), Inches(6.0), Inches(1.6),
         fill_color=BLUE_L, border_color=BLUE)
add_rich(slide, Inches(0.7), Inches(2.95), Inches(5.6), Inches(1.5), [
    ("实现语言决定性能天花板", 17, BLUE, True, PP_ALIGN.LEFT),
    ("", 3, DARK, False, PP_ALIGN.LEFT),
    ("Trie accept (55.7%) > Suffix (54.2%) — 算法质量更高", 13, DARK, False, PP_ALIGN.LEFT),
    ("但 speedup Suffix (2.04x) > Trie (1.95x) — C++ 胜出", 13, DARK, False, PP_ALIGN.LEFT),
    ("Python dict/DFS 开销 vs C++ 原生 O(1) 兄弟链", 12, BLUE, False, PP_ALIGN.LEFT),
])

# Insight 4
add_card(slide, Inches(6.8), Inches(2.9), Inches(6.0), Inches(1.6),
         fill_color=ORANGE_L, border_color=ORANGE)
add_rich(slide, Inches(7.0), Inches(2.95), Inches(5.6), Inches(1.5), [
    ("3-gram root 本质 = min_match_len", 17, ORANGE, True, PP_ALIGN.LEFT),
    ("", 3, DARK, False, PP_ALIGN.LEFT),
    ("Trie node_size=3: 需 3 token 匹配才进入查询", 13, DARK, False, PP_ALIGN.LEFT),
    ("Suffix min_match_len=3: 匹配长度 >= 3 才返回", 13, DARK, False, PP_ALIGN.LEFT),
    ("本质相同: 最低上下文长度保证预测质量", 12, ORANGE, False, PP_ALIGN.LEFT),
])

# 核心结论
add_card(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.4),
         fill_color=BG_CARD, border_color=BLUE, border_width=2)
add_rich(slide, Inches(0.8), Inches(4.85), Inches(11.7), Inches(2.3), [
    ("核心结论", 20, BLUE, True, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("技术路线: KMP (baseline) -> Hash (统计过滤) -> Trie (树结构+频率) -> Suffix (C++ 极致)", 14, DARK, False, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("C++ Suffix Tree 2.04x 是纯 N-gram 统计方案的实际上限", 15, RED, True, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("进一步突破 (如 2.5x+) 需要从 \"统计匹配\" 升级到 \"模型级 draft\" (如小模型蒸馏/Medusa)", 14, DARK, False, PP_ALIGN.LEFT),
])


# ══════════════════════════════════════════════════════════════
# Slide 14: 技术演进路线
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
title_bar(slide, "技术演进路线图", color=BLUE)

levels = [
    ("Level 1: 无状态匹配", "KMP",
     "每次从头扫描全序列, 无持久化, 无统计\nSpeedup: 1.70x",
     BLUE, BLUE_L, Inches(1.0)),
    ("Level 2: 统计过滤", "Hash",
     "频率表 + 置信度 + 多 n-gram 投票 + 请求局部叠加\nSpeedup: 1.86x  (+9%)",
     GREEN, GREEN_L, Inches(2.5)),
    ("Level 3: 树结构索引", "Trie",
     "三层 Trie + 输入/输出频率混合 + DFS 多分支展开\nSpeedup: 1.95x  (+5%)",
     PURPLE, PURPLE_L, Inches(4.0)),
    ("Level 4: C++ 极致", "Suffix Tree",
     "压缩后缀树 + 路径压缩 + 兄弟链 O(1) 最优 + 内建概率\nSpeedup: 2.04x  (+5%)",
     RED, RED_L, Inches(5.5)),
]

for title, name, desc, color, fill, y in levels:
    add_flow_box(slide, Inches(0.5), y, Inches(2.5), Inches(1.1),
                 title, border_color=color, fill_color=fill, font_size=12)
    add_flow_box(slide, Inches(3.3), y, Inches(1.8), Inches(1.1),
                 name, border_color=color, fill_color=color,
                 font_size=18, font_color=RGBColor(0xFF, 0xFF, 0xFF))
    add_card(slide, Inches(5.4), y, Inches(7.4), Inches(1.1),
             fill_color=fill, border_color=color)
    add_text(slide, Inches(5.6), y + Inches(0.1),
             Inches(7.0), Inches(0.9),
             desc, font_size=13, color=DARK)
    if y < Inches(5.5):
        add_down_arrow(slide, Inches(2.4), y + Inches(1.1),
                       Inches(0.3), Inches(0.3), color=color)


# ══════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════
output_path = "/home/hank/Agent/ngram/docs/ngram_speculative_decoding.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
